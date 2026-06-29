#!/usr/bin/env python3
"""Generate PlantMind handover Word doc and PowerPoint from consolidated blueprint."""

from pathlib import Path
from docx import Document
from docx.shared import Inches, Pt, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.table import WD_TABLE_ALIGNMENT
from pptx import Presentation
from pptx.util import Inches as PInches, Pt as PPt
from pptx.dml.color import RGBColor as PRGB
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

ROOT = Path(__file__).resolve().parent.parent
OUT = ROOT / "docs" / "deliverables"
OUT.mkdir(parents=True, exist_ok=True)

# Palette — Ocean Gradient (industrial trust)
NAVY = RGBColor(0x06, 0x5A, 0x82)
TEAL = RGBColor(0x1C, 0x72, 0x93)
MID = RGBColor(0x21, 0x29, 0x5C)
PNAVY = PRGB(0x06, 0x5A, 0x82)
PTEAL = PRGB(0x1C, 0x72, 0x93)
PMID = PRGB(0x21, 0x29, 0x5C)
PWHITE = PRGB(0xFF, 0xFF, 0xFF)
PLIGHT = PRGB(0xF0, 0xF4, 0xF8)


def add_heading(doc, text, level=1):
    h = doc.add_heading(text, level=level)
    for run in h.runs:
        run.font.color.rgb = NAVY if level == 1 else TEAL
    return h


def add_table(doc, headers, rows):
    t = doc.add_table(rows=1 + len(rows), cols=len(headers))
    t.style = "Table Grid"
    t.alignment = WD_TABLE_ALIGNMENT.CENTER
    for i, h in enumerate(headers):
        cell = t.rows[0].cells[i]
        cell.text = h
        for p in cell.paragraphs:
            for r in p.runs:
                r.bold = True
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            t.rows[ri + 1].cells[ci].text = str(val)
    doc.add_paragraph()


def build_word_doc():
    doc = Document()
    s = doc.sections[0]
    s.top_margin = s.bottom_margin = s.left_margin = s.right_margin = Inches(1)

    # Title page
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r = p.add_run("PlantMind × Götze Engine\n")
    r.bold = True
    r.font.size = Pt(28)
    r.font.color.rgb = NAVY
    r2 = p.add_run("\nComplete Project Blueprint\nConsolidated Specification v1.0")
    r2.font.size = Pt(16)
    r2.font.color.rgb = TEAL
    doc.add_paragraph()
    meta = doc.add_paragraph()
    meta.alignment = WD_ALIGN_PARAGRAPH.CENTER
    meta.add_run(
        "LTTS Global Engineering Intelligence Hackathon · 9 July 2026\n"
        "Owner: Sourav Dutta · LTTS\n"
        "Workspace: PlantMind-Live\n"
        "Date: 29 June 2026"
    )
    doc.add_page_break()

    # TOC note
    add_heading(doc, "Document Guide", 1)
    doc.add_paragraph(
        "This document consolidates PlantMind v1 (FORGE) and v2 (vault) into a single "
        "end-to-end blueprint. Audience routing:"
    )
    add_table(doc, ["Audience", "Primary Sections"], [
        ("CEO / Executive", "§1 Summary, §2 Why Win, §12 ROI, §16 Status"),
        ("Delivery Manager", "§10 Build Plan, §15 Governance, §16 Roadmap, §17 Team"),
        ("Technical Lead", "§4–9 Architecture, Data Flow, Agents, Models"),
        ("Peers / Engineers", "§5–8, §11 Stack, Appendix Contracts"),
        ("Judges / Handover", "§1–3, §9 Demo, §13 Research"),
    ])

    add_heading(doc, "§1 Executive Summary", 1)
    doc.add_paragraph(
        "PlantMind is a Physics-Informed Engineering Intelligence framework that converts "
        "industrial sensor data into one explainable, human-approved corrective action at "
        "the moment of maximum asset stress."
    )
    add_table(doc, ["Field", "Value"], [
        ("Product", "PlantMind × Götze Engine"),
        ("Tagline", "Predict the failure. Decide the fix. Prove it."),
        ("Strategic anchor", "LTTS × Databricks partnership (11 June 2026)"),
        ("Differentiator", "Decides optimal fix + proves with counterfactual + audit"),
        ("Team", "4 members · 24h OpenHack · 9 July 2026"),
        ("Status", "v1 FORGE built · v2 architecture locked · merge complete"),
    ])
    doc.add_paragraph(
        "The Götze analogy: In 2014, analytics told Germany's coach the single best "
        "substitution was Mario Götze — he scored the World Cup winner. PlantMind is "
        "that coach for a factory."
    )

    add_heading(doc, "§2 Problem & Market Validation", 1)
    add_table(doc, ["Metric", "Value"], [
        ("US unplanned downtime", "~$50B/year"),
        ("Fortune 500 downtime loss", "~$1.4T/year"),
        ("Plant cost per hour", "$50K–$500K"),
        ("Failures caught early", "Only 20–30%"),
        ("Industrial data analyzed", "<1%"),
    ])
    doc.add_paragraph(
        "Market gap: Every major vendor predicts failure. Almost none closes the loop to "
        "'here is the single best action, ranked, with reasoning, human-approved, and logged.'"
    )

    add_heading(doc, "§3 Consolidated Vision (Post-Merge)", 1)
    doc.add_paragraph("Three complementary views of one product:")
    for item in [
        "Layer 0 — 7 tool-agnostic interface contracts (LTTS IP)",
        "Layer 1 — Local hackathon build + Databricks production narrative",
        "Runtime — 5 specialist agents in LangGraph sequence",
        "v1 Reference — forge-v1: runnable MetaGPT pipeline with RED→GREEN proof",
    ]:
        doc.add_paragraph(item, style="List Bullet")

    add_heading(doc, "§4 System Architecture", 1)
    add_heading(doc, "Layer 0 — Interface Contracts", 2)
    add_table(doc, ["Contract", "Responsibility"], [
        ("IngestorInterface", "Sensor/maintenance ingest + validation"),
        ("FeatureStoreInterface", "Point-in-time features incl. physics"),
        ("AnomalyModelInterface", "Per-asset anomaly detection"),
        ("PhysicsModelInterface", "Health + RUL callable contract"),
        ("KnowledgeRetrieverInterface", "RAG over manuals/logs"),
        ("AgentOrchestratorInterface", "Agent trigger, routing, audit"),
        ("GovernanceInterface", "Lineage, policy, audit"),
        ("FeedbackLoopInterface", "Outcomes → recalibration"),
    ])

    add_heading(doc, "Layer 1 — Hackathon vs Production", 2)
    add_table(doc, ["Component", "Hackathon", "Production"], [
        ("Ingestor", "Python + SQLite/CSV", "Auto Loader + DLT"),
        ("Physics", "scipy Weibull", "Python wheel on cluster"),
        ("Orchestrator", "CrewAI + LangGraph", "Mosaic AI Agents"),
        ("Governance", "SQLite hash-chain audit", "Unity Catalog"),
        ("RAG", "ChromaDB + MiniLM", "Vector Search"),
    ])

    add_heading(doc, "§5 End-to-End Data Flow", 1)
    flow = (
        "Sensors → Bronze → Silver → Gold features → AssetHealthOracle (Weibull) → "
        "[TRIGGER] → GötzeEngine (IIS) → RootCauseAnalyst (RAG) → ExecutiveSummarizer → "
        "Dashboard → Human Approve → Immutable Audit → Feedback store"
    )
    doc.add_paragraph(flow)
    add_heading(doc, "Synthetic Demo Universe", 2)
    add_table(doc, ["Parameter", "Value"], [
        ("Assets", "30 (pump, compressor, motor, bearing, valve)"),
        ("Signals", "20 per asset"),
        ("Failure modes", "3 (gradual, sudden, intermittent)"),
        ("Cycles", "500 per asset"),
        ("Kaggle calibration", "CMAPSS, PRONOSTIA"),
    ])

    add_heading(doc, "§6 Agentic Workflows", 1)
    add_table(doc, ["#", "Agent", "Role", "Approval?"], [
        ("1", "DataSentinel", "Flag anomalies (Z-score + Mahalanobis)", "No"),
        ("2", "AssetHealthOracle", "Health 0–100 + RUL days + CI", "No"),
        ("3", "GötzeEngine ⭐", "IIS → ONE best action", "YES"),
        ("4", "RootCauseAnalyst", "RAG causal chain + citations", "No"),
        ("5", "ExecutiveSummarizer", "3-bullet leadership brief", "No"),
    ])
    doc.add_paragraph("Sequence: Sentinel → Oracle → Götze → RootCause → Summarizer → Human → Audit")

    add_heading(doc, "§7 Decision Engine — IIS (Canonical)", 1)
    doc.add_paragraph(
        "IIS(i) = 0.35·ΔP_failure + 0.25·ΔDowntimeCost + 0.20·Feasibility "
        "+ 0.15·HistoricalSuccess − 0.05·SafetyRiskDelta"
    )
    doc.add_paragraph(
        "v1 G-score (forge-v1): G = 0.40·ΔHealth + 0.25·(1−NormCost) + 0.20·(1−NormTime) "
        "+ 0.15·Safety — retained as runnable reference; IIS is canonical for product IP."
    )
    doc.add_paragraph("Proof stack: IIS ranking + RED→GREEN chart + audit hash-chain + human approval.")

    add_heading(doc, "§8 Model Routing", 1)
    add_table(doc, ["Task", "Model", "Rule"], [
        ("IIS / scoring", "Deterministic Python", "Never LLM"),
        ("Narrative", "Groq Llama 3.3 70B", "Fallback: templates"),
        ("Summaries", "Groq Llama 3.2 3B", "Fast labels"),
        ("Root cause", "DeepSeek R1 / Groq", "RAG-grounded only"),
        ("Embeddings", "all-MiniLM-L6-v2", "Local, free"),
        ("Health (canonical)", "scipy Weibull", "Always ships first"),
        ("RUL (v1 demo)", "RandomForest", "C-MAPSS fallback"),
        ("PINN", "PyTorch", "Stretch — freeze Hour 14"),
    ])
    doc.add_paragraph("Registry: ops/MODEL-REGISTRY.md · Adapters in src/agents/_llm.py")

    add_heading(doc, "§9 API Routing & Contracts", 1)
    routes = [
        "GET /api/v1/assets — list assets",
        "GET /api/v1/assets/{id}/health — AssetHealthReport",
        "POST /api/v1/assets/{id}/evaluate — trigger pipeline",
        "GET /api/v1/assets/{id}/decision — GötzeDecision",
        "POST /api/v1/decisions/{id}/approve — audit write",
        "GET /api/v1/audit — audit trail",
        "POST /api/v1/scenarios/inject — demo injector",
    ]
    for r in routes:
        doc.add_paragraph(r, style="List Bullet")
    doc.add_paragraph("Dashboard consumes JSON contracts only — never imports physics internals.")

    add_heading(doc, "§10 Demo Scenarios", 1)
    add_table(doc, ["Scenario", "Proves"], [
        ("A — Gradual pump wear ⭐", "One-best-action magic"),
        ("B — Sudden bearing impact", "Emergency stop urgency"),
        ("C — Intermittent valve", "Mahalanobis pattern depth"),
        ("D — Sensor dropout", "Data quality vs machine fault"),
    ])
    doc.add_paragraph("v1 demo today: streamlit run src/legacy/forge-v1/app.py")

    add_heading(doc, "§11 Technology Stack", 1)
    add_table(doc, ["Layer", "Technology"], [
        ("Backend", "Python 3.11, FastAPI, Pydantic v2"),
        ("Agents", "CrewAI + LangGraph"),
        ("LLM", "Groq Llama 3.3 70B / 3.2 3B"),
        ("Vector", "ChromaDB + MiniLM embeddings"),
        ("Physics", "scipy Weibull; PyTorch PINN (stretch)"),
        ("Dashboard", "Streamlit + Plotly"),
        ("Storage", "SQLite → Delta Lake"),
        ("Cloud", "Databricks DLT, Feature Store, Unity Catalog"),
    ])

    add_heading(doc, "§12 ROI Case", 1)
    doc.add_paragraph(
        "Value = (failures_prevented) × (avg_downtime_hours) × (cost_per_hour) − intervention_cost"
    )
    doc.add_paragraph("cost_per_hour: [ESTIMATE $50K–$500K/hr by sector]")

    add_heading(doc, "§13 Research & Competitive Position", 1)
    doc.add_paragraph("Completed: 35-pain register + 12+ competitor map (June 2026).")
    doc.add_paragraph(
        "Wedge: Open portable framework + physics IIS + human-governed audit + "
        "Databricks reference — vs predict-only incumbents."
    )

    add_heading(doc, "§14 Governance & IP", 1)
    for g in [
        "Non-autonomous — human approves all actions",
        "Explainable — reason + citations required",
        "Logged — immutable hash-chain audit",
        "Patent-candidate pending prior-art — closed decision loop is the wedge",
    ]:
        doc.add_paragraph(g, style="List Bullet")

    add_heading(doc, "§15 24-Hour Build Plan", 1)
    add_table(doc, ["Hours", "Milestone"], [
        ("0–4", "Scaffold, synthetic data, stubs"),
        ("4–8", "DataSentinel + AssetHealthOracle"),
        ("H8 ✓", "Health visible on dashboard"),
        ("8–12", "GötzeEngine + IIS"),
        ("12–16", "RootCause + ExecutiveSummarizer"),
        ("H16 ✓", "DEMO SAFE — full 5-agent flow"),
        ("16–20", "PINN stretch or harden"),
        ("H20", "Feature freeze"),
        ("20–23", "Rehearse 5+ times"),
        ("H23", "Tag v1.0-hackathon-submission"),
    ])

    add_heading(doc, "§16 Implementation Status", 1)
    add_table(doc, ["Component", "v1 forge", "v2 target"], [
        ("C-MAPSS ingest", "✅ Built", "✅"),
        ("Weibull multi-asset", "❌", "✅ Specified"),
        ("IIS scorer", "❌", "✅ Locked"),
        ("G-score", "✅ Built", "v1 alias"),
        ("5 agents", "❌", "✅ Specified"),
        ("RED→GREEN proof", "✅ Built", "✅"),
        ("Human approve", "❌", "✅ Specified"),
        ("FastAPI", "❌", "Scaffolded"),
        ("Streamlit", "✅ Built", "Migrate"),
    ])

    add_heading(doc, "§17 Team & Lanes", 1)
    add_table(doc, ["Lane", "Owner", "Deliverable"], [
        ("1 Backend", "Sourav", "Agents, API, orchestrator, audit"),
        ("2 Physics/ML", "Sourav", "Weibull, synthetic data, training"),
        ("3 Dashboard", "Member 2/3", "Streamlit, IIS panel, audit"),
        ("4 Databricks", "Sourav/team", "DLT, Feature Store, Unity"),
        ("5 Demo/Pitch", "Member 4/Sourav", "Script, ROI, Q&A"),
    ])

    add_heading(doc, "Appendix — Conflict Resolution Summary", 1)
    add_table(doc, ["Conflict", "Resolution"], [
        ("5 Layers vs 5 Agents", "Layers = pipeline; Agents = runtime — both valid"),
        ("G-score vs IIS", "IIS canonical; G-score = v1 implementation"),
        ("RF vs Weibull", "Weibull canonical; RF = C-MAPSS demo path"),
        ("MetaGPT vs LangGraph", "LangGraph canonical; MetaGPT = v1 POC"),
        ("No approval vs required", "Human approval required in v2"),
    ])
    doc.add_paragraph("Full matrix: docs/CONFLICT-RESOLUTION.md")
    doc.add_paragraph("Nothing deleted from source folders. 246+ files in PlantMind-Live.")

    out = OUT / "PlantMind_Complete_Project_Blueprint.docx"
    doc.save(str(out))
    print(f"Word: {out}")
    return out


def set_slide_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_title_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(slide, PNAVY)
    box = slide.shapes.add_textbox(PInches(0.5), PInches(2.2), PInches(9), PInches(1.5))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = PPt(40)
    p.font.bold = True
    p.font.color.rgb = PWHITE
    p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = subtitle
    p2.font.size = PPt(18)
    p2.font.color.rgb = PWHITE
    p2.alignment = PP_ALIGN.CENTER


def add_content_slide(prs, title, bullets, dark=False):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, PNAVY if dark else PLIGHT)
    tc = PWHITE if dark else PNAVY
    bc = PWHITE if dark else PMID
    # Title
    tb = slide.shapes.add_textbox(PInches(0.5), PInches(0.3), PInches(9), PInches(0.8))
    tfp = tb.text_frame.paragraphs[0]
    tfp.text = title
    tfp.font.size = PPt(32)
    tfp.font.bold = True
    tfp.font.color.rgb = tc
    # Accent bar
    bar = slide.shapes.add_shape(1, PInches(0.5), PInches(1.05), PInches(1.2), PInches(0.08))
    bar.fill.solid()
    bar.fill.fore_color.rgb = PTEAL
    bar.line.fill.background()
    # Body
    bb = slide.shapes.add_textbox(PInches(0.5), PInches(1.3), PInches(9), PInches(5.5))
    bf = bb.text_frame
    bf.word_wrap = True
    for i, b in enumerate(bullets):
        para = bf.paragraphs[0] if i == 0 else bf.add_paragraph()
        para.text = b
        para.font.size = PPt(16 if not dark else 18)
        para.font.color.rgb = bc
        para.space_after = PPt(10)
        para.level = 0


def add_stat_slide(prs, stats):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, PLIGHT)
    tb = slide.shapes.add_textbox(PInches(0.5), PInches(0.3), PInches(9), PInches(0.7))
    tb.text_frame.paragraphs[0].text = "The Problem — By the Numbers"
    tb.text_frame.paragraphs[0].font.size = PPt(32)
    tb.text_frame.paragraphs[0].font.bold = True
    tb.text_frame.paragraphs[0].font.color.rgb = PNAVY
    cols = len(stats)
    w = 8.5 / cols
    for i, (num, label) in enumerate(stats):
        left = 0.5 + i * w
        nb = slide.shapes.add_textbox(PInches(left), PInches(2), PInches(w - 0.2), PInches(1.2))
        np = nb.text_frame.paragraphs[0]
        np.text = num
        np.font.size = PPt(36)
        np.font.bold = True
        np.font.color.rgb = PTEAL
        np.alignment = PP_ALIGN.CENTER
        lb = slide.shapes.add_textbox(PInches(left), PInches(3.2), PInches(w - 0.2), PInches(1))
        lp = lb.text_frame.paragraphs[0]
        lp.text = label
        lp.font.size = PPt(12)
        lp.font.color.rgb = PNAVY
        lp.alignment = PP_ALIGN.CENTER


def build_ppt():
    prs = Presentation()
    prs.slide_width = PInches(10)
    prs.slide_height = PInches(7.5)

    add_title_slide(prs,
        "PlantMind × Götze Engine",
        "Engineering Intelligence · LTTS Hackathon · 9 July 2026\n"
        "Complete Project Handover · Consolidated Blueprint v1.0")

    add_content_slide(prs, "Executive Summary", [
        "Physics-Informed Engineering Intelligence for industrial assets",
        "At peak asset stress: scores every intervention → surfaces ONE best action",
        "Human approves · Everything logged · Counterfactual proof",
        "Strategic anchor: LTTS × Databricks partnership (June 2026)",
        "Tagline: Predict the failure. Decide the fix. Prove it.",
        "Status: v1 demo built · v2 architecture locked · merge complete",
    ])

    add_content_slide(prs, "The Götze Story", [
        "2014: Analytics told Germany's coach the single best substitution was Götze",
        "He came on and scored the World Cup winner",
        "PlantMind is that coach for a factory",
        "Everyone predicts failure — we decide the optimal fix and prove it works",
        "Memorable story + real partnership + working demo = what judges remember",
    ])

    add_stat_slide(prs, [
        ("$50B", "US unplanned\ndowntime / yr"),
        ("$1.4T", "Fortune 500\ndowntime loss"),
        ("<1%", "Industrial data\nactually analyzed"),
        ("20-30%", "Failures caught\nbefore occurrence"),
    ])

    add_content_slide(prs, "Consolidated Architecture", [
        "Layer 0: 7 tool-agnostic interface contracts (LTTS IP)",
        "Layer 1: Local hackathon build + Databricks production narrative",
        "Runtime: 5 agents — Sentinel → Oracle → Götze → RCA → Summarizer",
        "v1 Reference: forge-v1 MetaGPT pipeline with RED→GREEN proof (runnable)",
        "Conflict resolved: IIS canonical · Weibull canonical · LangGraph canonical",
    ])

    add_content_slide(prs, "End-to-End Data Flow", [
        "Sensors → Bronze → Silver → Gold features",
        "AssetHealthOracle: Weibull H(t) → health + RUL days + CI",
        "Trigger: health<40 OR rul<14 days OR critical anomaly",
        "GötzeEngine: IIS scores all interventions → ONE winner",
        "RootCauseAnalyst: RAG + citations → ExecutiveSummarizer: 3 bullets",
        "Dashboard → Human Approve → Immutable audit → Feedback",
    ])

    add_content_slide(prs, "The 5 Agents", [
        "1. DataSentinel — Z-score + Mahalanobis anomaly flags",
        "2. AssetHealthOracle — Weibull health 0-100 + RUL + 95% CI",
        "3. GötzeEngine ⭐ — IIS → one best action (HUMAN APPROVAL REQUIRED)",
        "4. RootCauseAnalyst — RAG over manuals with citations",
        "5. ExecutiveSummarizer — 3-bullet leadership brief + ROI",
        "Orchestration: LangGraph directed sequence with audit at every hop",
    ])

    add_content_slide(prs, "IIS — Intervention Impact Score", [
        "IIS = 0.35·ΔP_failure + 0.25·ΔDowntimeCost + 0.20·Feasibility",
        "      + 0.15·HistoricalSuccess − 0.05·SafetyRiskDelta",
        "All terms normalized [0,1] · Demo uses fixed weights",
        "SafetyRiskDelta above ceiling → hard veto regardless of score",
        "Proof: IIS breakdown + RED→GREEN chart + audit hash-chain",
        "v1 G-score retained in forge-v1 as runnable reference implementation",
    ])

    add_content_slide(prs, "Model Routing", [
        "Rule: AI does uncertain work · Deterministic rules decide",
        "IIS / G-score: Pure Python math — never LLM",
        "Narrative: Groq Llama 3.3 70B (fallback: templates)",
        "Root cause: DeepSeek R1 / Groq — RAG-grounded only",
        "Embeddings: all-MiniLM-L6-v2 local · Health: scipy Weibull",
        "Registry: ops/MODEL-REGISTRY.md · Adapters in src/agents/",
    ])

    add_content_slide(prs, "API & Contracts", [
        "FastAPI routes: /assets, /health, /evaluate, /decision, /approve, /audit",
        "Shared contracts (Pydantic): PhysicsModel, GötzeDecision, AuditRecord",
        "Dashboard reads JSON only — never imports physics internals",
        "Lane discipline: modules talk through src/contracts/ only",
        "Contract change = LOCKED_STATE vault update required",
    ])

    add_content_slide(prs, "Demo Scenarios", [
        "A ⭐ Gradual pump wear — hero one-best-action moment",
        "B Sudden bearing impact — emergency stop, high IIS gap",
        "C Intermittent valve — Mahalanobis catches hidden pattern",
        "D Sensor dropout — flags bad data, not bad machine",
        "v1 demo: streamlit run src/legacy/forge-v1/app.py",
        "5-min script: Hook → healthy plant → failure → approve → edge case → close",
    ])

    add_content_slide(prs, "LTTS × Databricks Alignment", [
        "Partnership announced 11 June 2026 — PlantMind is reference implementation",
        "Predictive Asset Reliability → AssetHealthOracle",
        "Energy & Emissions → IIS downtime-cost term",
        "OEE & Production → Downtime scoring",
        "Quality Intelligence → DataSentinel",
        "Sustainability Analytics → ExecutiveSummarizer ROI",
    ])

    add_content_slide(prs, "Governance & IP", [
        "Non-autonomous: human approves every action",
        "Explainable: reason + citations on every recommendation",
        "Logged: immutable hash-chain audit with full lineage",
        "IP wedge: closed loop physics → IIS → approved action → audit",
        "Framing: patent-candidate pending prior-art review",
        "Strongest claims: counterfactual proof + IIS scoring method",
    ])

    add_content_slide(prs, "24-Hour Build Plan", [
        "H0-4: Scaffold + synthetic data + agent stubs",
        "H4-8: DataSentinel + AssetHealthOracle",
        "H8 ✓ Health visible · H12: GötzeEngine + IIS",
        "H16 ✓ Full 5-agent flow — DEMO IS SAFE",
        "H20: Feature freeze · H23: Tag v1.0-hackathon-submission",
        "Golden rule: Not integrated by H16 → not in demo",
    ])

    add_content_slide(prs, "Implementation Status", [
        "✅ Built: C-MAPSS ingest, G-score, RED→GREEN, Streamlit, fleet view",
        "✅ Locked: IIS formula, 5 agents, contracts, Weibull params",
        "🔲 Pending: FastAPI, LangGraph orchestrator, RAG, human approve UI",
        "Strategy: Ship v1 demo if time-critical; migrate to v2 incrementally",
        "Workspace: PlantMind-Live (246+ files, nothing lost from merge)",
    ])

    add_content_slide(prs, "Team & Next Steps", [
        "Lane 1 Sourav: Backend, agents, API, audit",
        "Lane 2 Sourav: Weibull, synthetic data, ML training",
        "Lane 3 Members 2/3: Streamlit, IIS panel, audit views",
        "Lane 4: Databricks DLT, Feature Store, Unity Catalog",
        "Lane 5: Demo script, pitch, ROI, Q&A",
        "Start: 00-START-HERE.md → LOCKED_STATE → ROADMAP NOW",
    ], dark=True)

    add_title_slide(prs,
        "PlantMind × Götze Engine",
        "Predict the failure. Decide the fix. Prove it.\n"
        "Sourav Dutta · LTTS · PlantMind-Live · June 2026")

    out = OUT / "PlantMind_Complete_Handover_Deck.pptx"
    prs.save(str(out))
    print(f"PPT: {out}")
    return out


if __name__ == "__main__":
    build_word_doc()
    build_ppt()
    print("Done.")